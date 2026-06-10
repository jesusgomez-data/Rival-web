"""
RivalFit — Presentación de Ventas B2B
Genera el archivo PPTX con estilo de marca oficial.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Brand Tokens ────────────────────────────────────────────────────────────
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK        = RGBColor(0x0A, 0x0A, 0x0A)
CARD        = RGBColor(0x11, 0x11, 0x11)
MUTED_BG    = RGBColor(0x1F, 0x29, 0x37)
RED         = RGBColor(0xEF, 0x44, 0x44)
RED_DEEP    = RGBColor(0xDC, 0x26, 0x26)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_MUTED  = RGBColor(0x9C, 0xA3, 0xAF)
GREEN       = RGBColor(0x22, 0xC5, 0x5E)

W = Inches(13.33)   # Widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # totalmente en blanco


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=24, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bullet_block(slide, items, x, y, w, h,
                      font_size=18, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def slide_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def red_accent_bar(slide, top=True):
    """Thin red line at top or bottom."""
    y = Inches(0) if top else H - Inches(0.05)
    add_rect(slide, Inches(0), y, W, Inches(0.05), RED)


def footer(slide, url="rivalfit.app"):
    add_text(slide, url,
             Inches(0.4), H - Inches(0.45), Inches(4), Inches(0.35),
             font_size=9, color=GRAY_MUTED, align=PP_ALIGN.LEFT)
    add_text(slide, "© 2025 RivalFit",
             W - Inches(2.5), H - Inches(0.45), Inches(2.3), Inches(0.35),
             font_size=9, color=GRAY_MUTED, align=PP_ALIGN.RIGHT)


def slide_number_badge(slide, num):
    add_rect(slide, W - Inches(0.65), H - Inches(0.45), Inches(0.45), Inches(0.28), RED)
    add_text(slide, str(num),
             W - Inches(0.65), H - Inches(0.47), Inches(0.45), Inches(0.32),
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, BLACK)

# Diagonal red shape (decorative)
from pptx.enum.shapes import MSO_SHAPE_TYPE
# Large dark card background
add_rect(s, Inches(0), Inches(0), W, H, BLACK)

# Right side red gradient block
add_rect(s, W - Inches(4.2), Inches(0), Inches(4.2), H, RGBColor(0x1A, 0x04, 0x04))
add_rect(s, W - Inches(0.15), Inches(0), Inches(0.15), H, RED)

# Top red bar
add_rect(s, Inches(0), Inches(0), W, Inches(0.08), RED)

# Tag
add_rect(s, Inches(0.6), Inches(1.6), Inches(1.4), Inches(0.32), RED)
add_text(s, "B2B  SALES",
         Inches(0.6), Inches(1.58), Inches(1.5), Inches(0.36),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Main title
add_text(s, "RIVALFIT",
         Inches(0.55), Inches(2.05), Inches(8), Inches(1.8),
         font_size=96, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

# Red underline
add_rect(s, Inches(0.55), Inches(3.6), Inches(5.5), Inches(0.07), RED)

# Tagline
add_text(s, "La plataforma todo en uno para centros deportivos,\nboxes CrossFit y gimnasios modernos.",
         Inches(0.55), Inches(3.8), Inches(7.8), Inches(1.2),
         font_size=20, color=GRAY, align=PP_ALIGN.LEFT)

# Sub tagline
add_text(s, "ENTRENA · COMPITE · DOMINA",
         Inches(0.55), Inches(5.1), Inches(7), Inches(0.5),
         font_size=13, bold=True, italic=True,
         color=RGBColor(0xEF, 0x44, 0x44), align=PP_ALIGN.LEFT)

# URL
add_text(s, "rivalfit.app",
         Inches(0.55), H - Inches(0.7), Inches(3), Inches(0.4),
         font_size=12, color=GRAY_MUTED, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — El Problema
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 2)
footer(s)

# Section label
add_rect(s, Inches(0.6), Inches(0.55), Inches(1.4), Inches(0.28), RED)
add_text(s, "EL PROBLEMA",
         Inches(0.6), Inches(0.53), Inches(1.6), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "¿Cuánto tiempo pierdes cada día?",
         Inches(0.6), Inches(1.0), Inches(10), Inches(0.7),
         font_size=36, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_text(s, "La mayoría de centros siguen gestionando todo de forma manual y fragmentada.",
         Inches(0.6), Inches(1.75), Inches(9), Inches(0.55),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

# Two columns of problems
problems_left = [
    "❌  WhatsApp para comunicar todo",
    "❌  Excel para controlar atletas",
    "❌  Cobros manuales y tardíos",
    "❌  Sin seguimiento del rendimiento",
]
problems_right = [
    "❌  Mala comunicación interna",
    "❌  Baja retención de atletas",
    "❌  Demasiado trabajo administrativo",
    "❌  Imagen poco profesional",
]

add_bullet_block(s, problems_left,
                 Inches(0.6), Inches(2.5), Inches(5.8), Inches(3.5),
                 font_size=17, color=GRAY, spacing=Pt(10))
add_bullet_block(s, problems_right,
                 Inches(6.8), Inches(2.5), Inches(5.8), Inches(3.5),
                 font_size=17, color=GRAY, spacing=Pt(10))

# Bottom stat bar
add_rect(s, Inches(0.6), Inches(6.4), Inches(11.8), Inches(0.7), RGBColor(0x1A, 0x04, 0x04))
add_rect(s, Inches(0.6), Inches(6.4), Inches(0.06), Inches(0.7), RED)
add_text(s, "El 73% de los gimnasios pierde miembros por falta de seguimiento y comunicación.",
         Inches(0.85), Inches(6.42), Inches(11), Inches(0.65),
         font_size=12, italic=True, color=GRAY_MUTED, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — La Solución
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 3)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(1.4), Inches(0.28), RED)
add_text(s, "LA SOLUCIÓN",
         Inches(0.6), Inches(0.53), Inches(1.6), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Todo en una sola plataforma.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
         font_size=40, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_text(s, "RivalFit centraliza toda la gestión de tu centro deportivo.",
         Inches(0.6), Inches(1.75), Inches(9), Inches(0.5),
         font_size=16, color=GRAY, align=PP_ALIGN.LEFT)

# Dashboard placeholder box
add_rect(s, Inches(0.6), Inches(2.4), Inches(7.8), Inches(3.6), CARD)
add_rect(s, Inches(0.6), Inches(2.4), Inches(7.8), Inches(0.04), RED)
add_text(s, "[ Dashboard RivalFit ]",
         Inches(0.6), Inches(3.7), Inches(7.8), Inches(1.0),
         font_size=16, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

# Benefits column
benefits = [
    "✅  Gestión completa del centro",
    "✅  Automatización de procesos",
    "✅  Control total en tiempo real",
    "✅  Escalable para cualquier tamaño",
    "✅  App para atletas incluida",
    "✅  Soporte y onboarding dedicado",
]
add_bullet_block(s, benefits,
                 Inches(9.0), Inches(2.4), Inches(3.8), Inches(3.8),
                 font_size=14, color=GRAY, spacing=Pt(8))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dashboard Profesional
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 4)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(2.1), Inches(0.28), RED)
add_text(s, "FUNCIONALIDAD 01",
         Inches(0.6), Inches(0.53), Inches(2.3), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Dashboard Profesional",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_rect(s, Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.04), RED)

# Image placeholder
add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(4.2), CARD)
add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(0.04), RED)
add_text(s, "[ Imagen Dashboard ]",
         Inches(0.6), Inches(3.7), Inches(8.2), Inches(1.0),
         font_size=15, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

# Description
add_text(s, "Controla ingresos, atletas, asistencia y operaciones desde un único panel centralizado.",
         Inches(9.15), Inches(2.1), Inches(3.8), Inches(1.2),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

details = [
    "→  Ingresos en tiempo real",
    "→  Asistencia diaria",
    "→  Nuevos atletas",
    "→  Alertas y notificaciones",
    "→  Acceso multisede",
]
add_bullet_block(s, details,
                 Inches(9.15), Inches(3.5), Inches(3.8), Inches(3.0),
                 font_size=13, color=GRAY_MUTED, spacing=Pt(7))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Agenda Inteligente
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 5)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(2.1), Inches(0.28), RED)
add_text(s, "FUNCIONALIDAD 02",
         Inches(0.6), Inches(0.53), Inches(2.3), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Agenda Inteligente",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_rect(s, Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.04), RED)

add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(4.2), CARD)
add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(0.04), RED)
add_text(s, "[ Imagen Agenda / Calendario ]",
         Inches(0.6), Inches(3.7), Inches(8.2), Inches(1.0),
         font_size=15, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Organiza clases, coaches y aforo en segundos.",
         Inches(9.15), Inches(2.1), Inches(3.8), Inches(0.9),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

details = [
    "→  Vista semanal por colores",
    "→  Gestión de coaches",
    "→  Control de aforo",
    "→  Google Calendar sync",
    "→  Reservas de atletas",
    "→  Clases recurrentes",
]
add_bullet_block(s, details,
                 Inches(9.15), Inches(3.1), Inches(3.8), Inches(3.4),
                 font_size=13, color=GRAY_MUTED, spacing=Pt(7))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Base de Datos de Atletas
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 6)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(2.1), Inches(0.28), RED)
add_text(s, "FUNCIONALIDAD 03",
         Inches(0.6), Inches(0.53), Inches(2.3), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Base de Datos de Atletas",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_rect(s, Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.04), RED)

add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(4.2), CARD)
add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(0.04), RED)
add_text(s, "[ Imagen Base de Datos / Miembros ]",
         Inches(0.6), Inches(3.7), Inches(8.2), Inches(1.0),
         font_size=15, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Gestiona altas, bajas, membresías, pagos e historial completo de cada atleta.",
         Inches(9.15), Inches(2.1), Inches(3.8), Inches(1.1),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

details = [
    "→  Altas y bajas digitales",
    "→  Historial de pagos",
    "→  Membresías y bonos",
    "→  Documentos y contratos",
    "→  Ficha médica y lesiones",
    "→  Gestión de leads/trials",
]
add_bullet_block(s, details,
                 Inches(9.15), Inches(3.3), Inches(3.8), Inches(3.2),
                 font_size=13, color=GRAY_MUTED, spacing=Pt(7))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Check-In Automático
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 7)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(2.1), Inches(0.28), RED)
add_text(s, "FUNCIONALIDAD 04",
         Inches(0.6), Inches(0.53), Inches(2.3), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Check-In Automático",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_rect(s, Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.04), RED)

add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(4.2), CARD)
add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(0.04), RED)
add_text(s, "[ Imagen Check-In / Asistencia ]",
         Inches(0.6), Inches(3.7), Inches(8.2), Inches(1.0),
         font_size=15, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Control de acceso y asistencia sin procesos manuales.",
         Inches(9.15), Inches(2.1), Inches(3.8), Inches(0.9),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

details = [
    "→  Check-in desde app",
    "→  QR por clase",
    "→  Historial de asistencia",
    "→  Alertas de inactividad",
    "→  Registro automático",
]
add_bullet_block(s, details,
                 Inches(9.15), Inches(3.1), Inches(3.8), Inches(3.4),
                 font_size=13, color=GRAY_MUTED, spacing=Pt(7))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Gestión de WODs
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 8)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(2.1), Inches(0.28), RED)
add_text(s, "FUNCIONALIDAD 05",
         Inches(0.6), Inches(0.53), Inches(2.3), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Gestión de WODs",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_rect(s, Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.04), RED)

add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(4.2), CARD)
add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(0.04), RED)
add_text(s, "[ Imagen WOD / Entrenamiento del día ]",
         Inches(0.6), Inches(3.7), Inches(8.2), Inches(1.0),
         font_size=15, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Publica entrenamientos diarios para toda tu comunidad.",
         Inches(9.15), Inches(2.1), Inches(3.8), Inches(0.9),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

details = [
    "→  WOD diario publicado",
    "→  Leaderboard por clase",
    "→  Resultados de atletas",
    "→  Generador IA de WODs",
    "→  Historial completo",
    "→  Notificaciones push",
]
add_bullet_block(s, details,
                 Inches(9.15), Inches(3.1), Inches(3.8), Inches(3.4),
                 font_size=13, color=GRAY_MUTED, spacing=Pt(7))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Gestión de Equipo
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 9)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(2.1), Inches(0.28), RED)
add_text(s, "FUNCIONALIDAD 06",
         Inches(0.6), Inches(0.53), Inches(2.3), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Gestión de Equipo",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_rect(s, Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.04), RED)

add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(4.2), CARD)
add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(0.04), RED)
add_text(s, "[ Imagen Gestión de Equipo / Coaches ]",
         Inches(0.6), Inches(3.7), Inches(8.2), Inches(1.0),
         font_size=15, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Controla coaches, permisos y horarios desde una sola plataforma.",
         Inches(9.15), Inches(2.1), Inches(3.8), Inches(0.9),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

details = [
    "→  Roles: Owner / Head Coach",
    "→  Permisos por rol",
    "→  Asignación de clases",
    "→  Rendimiento del equipo",
    "→  Comunicación interna",
]
add_bullet_block(s, details,
                 Inches(9.15), Inches(3.1), Inches(3.8), Inches(3.4),
                 font_size=13, color=GRAY_MUTED, spacing=Pt(7))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Coach IA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 10)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(2.1), Inches(0.28), RED)
add_text(s, "FUNCIONALIDAD 07",
         Inches(0.6), Inches(0.53), Inches(2.3), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Coach IA Integrado",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_rect(s, Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.04), RED)

add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(4.2), CARD)
add_rect(s, Inches(0.6), Inches(2.1), Inches(8.2), Inches(0.04), RED)
add_text(s, "[ Imagen Coach IA / Asistente ]",
         Inches(0.6), Inches(3.7), Inches(8.2), Inches(1.0),
         font_size=15, color=GRAY_MUTED, align=PP_ALIGN.CENTER)

add_text(s, "Asistente deportivo disponible 24/7 para todos tus atletas.",
         Inches(9.15), Inches(2.1), Inches(3.8), Inches(0.9),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

details = [
    "→  Planes de entrenamiento IA",
    "→  Nutrición personalizada",
    "→  Análisis de rendimiento",
    "→  Respuestas instantáneas",
    "→  Disponible en la app",
    "→  Diferenciación premium",
]
add_bullet_block(s, details,
                 Inches(9.15), Inches(3.1), Inches(3.8), Inches(3.4),
                 font_size=13, color=GRAY_MUTED, spacing=Pt(7))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Resultados
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 11)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(2.2), Inches(0.28), RED)
add_text(s, "LO QUE CONSIGUES",
         Inches(0.6), Inches(0.53), Inches(2.4), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Tu centro. Transformado.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=42, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_text(s, "Con RivalFit, los centros reportan mejoras en todos los indicadores clave.",
         Inches(0.6), Inches(1.8), Inches(9.5), Inches(0.5),
         font_size=15, color=GRAY, align=PP_ALIGN.LEFT)

# 2-column results grid
col1 = [
    "✅  Menos tiempo administrativo",
    "✅  Gestión centralizada",
    "✅  Mejor experiencia para atletas",
    "✅  Mayor retención de miembros",
]
col2 = [
    "✅  Comunicación interna fluida",
    "✅  Comunidad propia en la app",
    "✅  Automatización de procesos",
    "✅  Imagen premium para tu marca",
]

add_bullet_block(s, col1, Inches(0.6), Inches(2.5), Inches(5.8), Inches(3.6),
                 font_size=17, color=GRAY, spacing=Pt(14))
add_bullet_block(s, col2, Inches(6.8), Inches(2.5), Inches(5.8), Inches(3.6),
                 font_size=17, color=GRAY, spacing=Pt(14))

# Stat boxes
stats = [
    ("-40%", "Tiempo\nadministrativo"),
    ("+35%", "Retención\nde atletas"),
    ("24/7", "Coach IA\ndisponible"),
    ("1 sola\napp", "Gestión\ncompleta"),
]
box_w = Inches(2.6)
box_h = Inches(0.9)
start_x = Inches(0.6)
y_stat = H - Inches(1.55)
gap = Inches(0.25)

for i, (num, label) in enumerate(stats):
    bx = start_x + i * (box_w + gap)
    add_rect(s, bx, y_stat, box_w, box_h, RGBColor(0x1A, 0x04, 0x04))
    add_rect(s, bx, y_stat, Inches(0.06), box_h, RED)
    add_text(s, num, bx + Inches(0.15), y_stat, Inches(1.1), box_h,
             font_size=18, bold=True, color=RED, align=PP_ALIGN.LEFT, font_name="Arial Black")
    add_text(s, label, bx + Inches(1.3), y_stat, Inches(1.2), box_h,
             font_size=10, color=GRAY_MUTED, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Planes
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK)
red_accent_bar(s, top=True)
slide_number_badge(s, 12)
footer(s)

add_rect(s, Inches(0.6), Inches(0.55), Inches(1.0), Inches(0.28), RED)
add_text(s, "PLANES",
         Inches(0.6), Inches(0.53), Inches(1.2), Inches(0.32),
         font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Elige el plan para tu centro.",
         Inches(0.6), Inches(1.0), Inches(12), Inches(0.75),
         font_size=38, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_rect(s, Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.04), RED)

plans = [
    ("STARTER",    "€49 / mes", "1 sede\nHasta 50 atletas\nAgenda + Check-in\nWODs + Leaderboard\nApp para atletas", False),
    ("PRO",        "€99 / mes", "Hasta 3 sedes\nAtletas ilimitados\nTodo Starter +\nCoach IA\nAnalíticas avanzadas", True),
    ("ENTERPRISE", "A medida",  "Cadenas y franquicias\nSedes ilimitadas\nAPI y personalizaciones\nSoporte prioritario\nOnboarding dedicado", False),
]

card_w  = Inches(3.8)
card_h  = Inches(4.5)
card_y  = Inches(2.1)
card_gap = Inches(0.35)

for i, (name, price, features, highlighted) in enumerate(plans):
    cx = Inches(0.6) + i * (card_w + card_gap)
    bg = RGBColor(0x1A, 0x04, 0x04) if highlighted else CARD
    add_rect(s, cx, card_y, card_w, card_h, bg)
    # Top accent
    top_color = RED if highlighted else RGBColor(0x33, 0x33, 0x33)
    add_rect(s, cx, card_y, card_w, Inches(0.06), top_color)

    add_text(s, name, cx + Inches(0.2), card_y + Inches(0.15), card_w - Inches(0.3), Inches(0.45),
             font_size=18, bold=True, italic=True, color=WHITE if highlighted else GRAY,
             align=PP_ALIGN.LEFT, font_name="Arial Black")
    add_text(s, price, cx + Inches(0.2), card_y + Inches(0.65), card_w - Inches(0.3), Inches(0.55),
             font_size=26, bold=True, color=RED if highlighted else WHITE,
             align=PP_ALIGN.LEFT, font_name="Arial Black")
    add_rect(s, cx + Inches(0.2), card_y + Inches(1.25), card_w - Inches(0.4), Inches(0.02),
             RGBColor(0x33, 0x33, 0x33))

    feat_items = ["→  " + f for f in features.split('\n')]
    add_bullet_block(s, feat_items,
                     cx + Inches(0.2), card_y + Inches(1.4), card_w - Inches(0.3), Inches(2.8),
                     font_size=12, color=GRAY, spacing=Pt(8))

    btn_color = RED if highlighted else RGBColor(0x2A, 0x2A, 0x2A)
    add_rect(s, cx + Inches(0.2), card_y + card_h - Inches(0.7),
             card_w - Inches(0.4), Inches(0.48), btn_color)
    add_text(s, "Solicita una demo",
             cx + Inches(0.2), card_y + card_h - Inches(0.72),
             card_w - Inches(0.4), Inches(0.52),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Cierre / CTA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, BLACK)

# Full red right strip
add_rect(s, W - Inches(0.15), Inches(0), Inches(0.15), H, RED)

# Top bar
add_rect(s, Inches(0), Inches(0), W, Inches(0.08), RED)

# Center quote block
add_rect(s, Inches(0.6), Inches(1.3), Inches(11.8), Inches(4.5), RGBColor(0x0D, 0x0D, 0x0D))
add_rect(s, Inches(0.6), Inches(1.3), Inches(0.1), Inches(4.5), RED)

add_text(s, "¿Listo para transformar\ntu centro deportivo?",
         Inches(1.0), Inches(1.5), Inches(11), Inches(1.8),
         font_size=46, bold=True, italic=True, color=WHITE, align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_text(s, "Agenda una demo gratuita y personalizada. Sin compromisos.",
         Inches(1.0), Inches(3.4), Inches(9), Inches(0.55),
         font_size=17, color=GRAY, align=PP_ALIGN.LEFT)

# CTA buttons area
add_rect(s, Inches(1.0), Inches(4.15), Inches(3.2), Inches(0.55), RED)
add_text(s, "📞  Agenda tu Demo",
         Inches(1.0), Inches(4.13), Inches(3.2), Inches(0.59),
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, Inches(4.5), Inches(4.15), Inches(3.2), Inches(0.55), RGBColor(0x22, 0x22, 0x22))
add_text(s, "🌐  rivalfit.app",
         Inches(4.5), Inches(4.13), Inches(3.2), Inches(0.59),
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "RIVALFIT",
         Inches(1.0), Inches(5.3), Inches(11), Inches(1.6),
         font_size=80, bold=True, italic=True,
         color=RGBColor(0x1A, 0x1A, 0x1A), align=PP_ALIGN.LEFT,
         font_name="Arial Black")

add_text(s, "ENTRENA · COMPITE · DOMINA",
         Inches(1.0), H - Inches(0.6), Inches(8), Inches(0.4),
         font_size=11, bold=True, italic=True, color=RED, align=PP_ALIGN.LEFT)

# ── Save ─────────────────────────────────────────────────────────────────────
out_dir  = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "RivalFit_Presentacion_Ventas.pptx")
prs.save(out_path)
print("Saved: " + out_path)
